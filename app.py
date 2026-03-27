import streamlit as st
import pandas as pd
import google.generativeai as genai
from pptx import Presentation
from io import StringIO

st.set_page_config(page_title="AdStrategix AI", layout="wide")

# ================= LOGIN FIXED =================
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if not st.session_state.logged_in:
    st.title("🔐 Login")

    user = st.text_input("Username")
    pwd = st.text_input("Password", type="password")

    if st.button("Login"):
        if user == "admin" and pwd == "1234":
            st.session_state.logged_in = True
            st.success("Login successful")
            st.rerun()  # FIX
        else:
            st.error("Invalid credentials")

    st.stop()

# ================= GEMINI =================
api_key = st.sidebar.text_input("Gemini API Key", type="password")

if api_key:
    genai.configure(api_key=api_key)
    st.sidebar.success("✅ API Key Loaded")

    # Test button
    if st.sidebar.button("Test API"):
        try:
            model = genai.GenerativeModel("gemini-pro")
            res = model.generate_content("Say hello")
            st.sidebar.success("✅ API Working")
        except Exception as e:
            st.sidebar.error(f"❌ Error: {str(e)}")

# ================= FUNCTIONS =================

def detect_platform(df):
    cols = [c.lower() for c in df.columns]

    if any("ad group" in c for c in cols):
        return "Google Ads"
    elif any("adset" in c for c in cols):
        return "Meta Ads"
    elif any("line item" in c for c in cols):
        return "DV360 / Programmatic"
    return "Unknown Platform"

def create_ppt(brand, industry, results):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"{brand} Media Plan"
    slide.placeholders[1].text = f"Industry: {industry}"

    for r in results:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = r["Platform"]
        slide.placeholders[1].text = f"Conversions: {r['Conversions']} | CPA: {r['CPA']}"

    prs.save("media_plan.pptx")

# ================= UI =================
st.title("🚀 AdStrategix AI - Paid Media Intelligence")

tab1, tab2, tab3, tab4 = st.tabs([
    "📊 Planner",
    "📈 Analyzer",
    "🧾 Ad Preview",
    "💾 Saved Plans"
])

# =====================================================
# TAB 1 — PLANNER
# =====================================================
with tab1:

    st.header("📊 Media Planner")

    brand = st.text_input("Brand Name")
    industry = st.selectbox("Industry", ["Real Estate", "SaaS", "Ecommerce"])
    country = st.selectbox("Country", ["India", "USA", "UK"])
    objective = st.selectbox("Objective", ["Leads", "Sales"])
    budget = st.number_input("Budget", min_value=1000)

    benchmarks = {
        "India": {"cpm": 250, "ctr": 1.2, "cvr": 3},
        "USA": {"cpm": 20, "ctr": 2, "cvr": 4},
        "UK": {"cpm": 18, "ctr": 1.8, "cvr": 3.5}
    }

    if st.button("Generate Plan"):

        data = benchmarks[country]
        results = []

        platforms = {
            "Google Ads": 0.4,
            "Meta Ads": 0.35,
            "LinkedIn": 0.15,
            "Programmatic": 0.1
        }

        st.subheader("📊 KPI Forecast")

        for p, split in platforms.items():
            sub_budget = budget * split

            impressions = (sub_budget / data["cpm"]) * 1000
            clicks = impressions * (data["ctr"] / 100)
            conversions = clicks * (data["cvr"] / 100)
            cpa = sub_budget / conversions if conversions else 0

            st.metric(p, f"{int(conversions)} conversions | CPA {round(cpa,2)}")

            results.append({
                "Platform": p,
                "Conversions": int(conversions),
                "CPA": round(cpa, 2)
            })

        # AI Persona
        st.subheader("👤 AI Buyer Persona")
        st.write(generate_ai(f"Create buyer persona for {industry} in {country}"))

        # AI Ads
        st.subheader("📝 AI Ad Copies")
        st.write(generate_ai(f"Generate ad copies for {brand} in {industry}"))

        # Keywords
        st.subheader("🔍 Keyword Strategy")

        kw_prompt = f"""
        Generate Google Ads keywords CSV with columns:
        Keyword, Type, Match Type

        For {industry} in {country}
        """

        kw = generate_ai(kw_prompt)
        st.code(kw)

        st.download_button("Download Keywords CSV", kw, "keywords.csv")

        # Budget Optimizer
        st.subheader("🧠 Budget Optimizer")
        st.write(generate_ai(f"Optimize budget for {industry} with {budget}"))

        # Save plan
        if "plans" not in st.session_state:
            st.session_state.plans = []

        st.session_state.plans.append({
            "brand": brand,
            "results": results
        })

        # PPT Export
        if st.button("Download PPT"):
            create_ppt(brand, industry, results)
            with open("media_plan.pptx", "rb") as f:
                st.download_button("Download Media Plan", f, "media_plan.pptx")

# =====================================================
# TAB 2 — ANALYZER
# =====================================================
with tab2:

    st.header("📈 Campaign Analyzer")

    file = st.file_uploader("Upload CSV", type=["csv"])

    if file:
        df = pd.read_csv(file)

        st.dataframe(df.head())

        platform = detect_platform(df)
        st.success(f"Detected Platform: {platform}")

        st.subheader("📊 Charts")
        st.bar_chart(df.select_dtypes(include='number'))

        # Creative Insights
        st.subheader("🎨 Creative Insights")

        for col in df.columns:
            if "ctr" in col.lower():
                top = df.sort_values(col, ascending=False).head(5)
                st.write("Top Performers")
                st.dataframe(top)

        # Custom AI Prompt
        st.subheader("💬 Custom AI Analysis")

        user_prompt = st.text_area("Ask anything about report")

        if st.button("Run AI Analysis"):
            sample = df.head(20).to_string()
            final_prompt = f"{user_prompt}\n\nData:\n{sample}"
            st.write(generate_ai(final_prompt))

# =====================================================
# TAB 3 — AD PREVIEW
# =====================================================
with tab3:

    st.header("🧾 Ad Tag Preview")

    code = st.text_area("Paste Ad Code")

    if st.button("Preview Ad"):
        st.components.v1.html(code, height=400)

# =====================================================
# TAB 4 — SAVED
# =====================================================
with tab4:

    st.header("💾 Saved Plans")

    if "plans" in st.session_state and st.session_state.plans:

        for plan in st.session_state.plans:
            st.subheader(plan["brand"])
            st.dataframe(pd.DataFrame(plan["results"]))

    else:
        st.info("No saved plans yet")
